import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image


class FileProcessor:
    def extract_content(self, file_path: str, file_type: str):
        """
        Extracts text from documents or prepares image paths for the model.
        """
        content = ""
        result = {"type": "text", "content": ""}
        try:
            # 1. Handle PDF
            if "pdf" in file_type:
                reader = PdfReader(file_path)
                for page in reader.pages:
                    text = page.extract_text()
                    if text: content += text + "\n"
                result["content"] = content
            # 2. Handle Word Documents
            elif "word" in file_type or "doc" in file_type:
                doc = Document(file_path)
                for para in doc.paragraphs:
                    content += para.text + "\n"
                result["content"] = content
            # 3. Handle PowerPoint
            elif "presentation" in file_type or "pptx" in file_type:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                result["content"] = content
            # 4. Handle Images 
            elif "image" in file_type:
                # For images, we don't extract text manually.
                # We return the path so Gemini can "see" it.
                result["type"] = "image"
                result["path"] = file_path
            # 5. Handle Plain Text
            elif "text" in file_type:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                result["content"] = content
            return result
        except Exception as e:
            return {"error": str(e)}